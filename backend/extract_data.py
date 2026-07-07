import os
import magic
import fitz
import pandas as pd
import docx
import pptx
import cv2
import pytesseract
import io
import numpy as np

def process_plain_text(file_stream):
    return file_stream.read().decode('utf-8')

def process_pdf(file_stream):
    File_Data = ""
    doc = fitz.open(stream=file_stream.read(), filetype="pdf")
    for page in doc:
        File_Data += page.get_text()
    return File_Data

def process_docx(file_stream):
    File_Data = ""
    file_bytes = io.BytesIO(file_stream.read())
    doc = docx.Document(file_bytes)
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            File_Data += paragraph.text.strip() + "\n"
    return File_Data

def process_ppt(file_stream):
    File_Data = ""
    file_bytes = io.BytesIO(file_stream.read())
    presentation = pptx.Presentation(file_bytes)
    for slide_num, slide in enumerate(presentation.slides, start=1):
        File_Data += f"\nSlide Number - {slide_num}\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                File_Data += shape.text.strip() + "\n"
    return File_Data

def process_image(file_stream):
    file_bytes = np.frombuffer(file_stream.read(), np.uint8)
    img = cv2.imdecode(file_bytes, cv2.IMREAD_COLOR)

    if img is None:
        return ""

    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    gray = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY)[1]

    text = pytesseract.image_to_string(gray, config="--psm 6")
    return text

def process_excel(file_stream):
    df = pd.read_excel(file_stream)
    return df.to_csv(index=False)


def extract_Data(uploaded_file):

    header_bytes = uploaded_file.read(2048)
    mime_type = magic.from_buffer(header_bytes, mime = True).lower()

    uploaded_file.seek(0) 

    SimpleTextExtensions = ["plain", "html", "xml", "json", "txt"]
    ImageExtensions = ["jpg", "jpeg", "png"]

    if "pdf" in mime_type:
        data = process_pdf(uploaded_file)

    elif "presentation" in mime_type:
        data = process_ppt(uploaded_file)
    
    elif "sheet" in mime_type or "excel" in mime_type:
        data = process_excel(uploaded_file)

    elif "document" in mime_type or "word" in mime_type:
        data = process_docx(uploaded_file)

    elif any(t in mime_type for t in SimpleTextExtensions):
        data = process_plain_text(uploaded_file)
    
    elif any(t in mime_type for t in ImageExtensions):
        data = process_image(uploaded_file)
    
    else:
        raise Exception("Invalid file type. That file format is not allowed.")

    return data
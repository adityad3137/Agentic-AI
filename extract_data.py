import os
import magic
import fitz
import pandas as pd
import docx
import pptx
import cv2
import pytesseract

def process_plain_text(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        File_Data = f.read()

    return File_Data

def process_pdf(file_path):
    File_Data = ""

    doc = fitz.open(file_path)
    for page in doc:
        File_Data = File_Data + page.get_text()

    return File_Data

def process_docx(file_path):
    File_Data = ""
    doc = docx.Document(file_path)

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            File_Data = File_Data + paragraph.text.strip() + "\n"
    
    return File_Data

def process_ppt(file_path):
    File_Data = ""
    presentation = pptx.Presentation(file_path)

    for slide_num, slide in enumerate(presentation.slides, start = 1):
        File_Data = File_Data + f"\nSlide Number - {slide_num}\n"

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                File_Data = File_Data + shape.text.strip() + "\n"
    
    return File_Data

def process_image(file_path):

    img = cv2.imread(file_path)

    if img is None:
        print("Failed to load image")
        return ""

    gray = cv2.cvtColor(
        img,
        cv2.COLOR_BGR2GRAY
    )

    gray = cv2.threshold(
        gray,
        150,
        255,
        cv2.THRESH_BINARY
    )[1]

    cv2.imwrite("debug.png", gray)

    text = pytesseract.image_to_string(
        gray,
        config="--psm 6"
    )

    return text


def process_excel(file_path):
    df = pd.read_excel(file_path)
    File_Data = df.to_csv(index = False)

    return File_Data


def extract_Data(file_path):

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"The file path does not exist: {file_path}")

    print(f"File found! Processing : {file_path}")

    mime_type = magic.from_file(file_path, mime = True).lower()

    SimpleTextExtensions = ["plain", "html", "xml", "json", "txt"]
    ImageExtensions = ["jpg", "jpeg", "png"]

    if "pdf" in mime_type:
        data = process_pdf(file_path)

    elif "presentation" in mime_type:
        data = process_ppt(file_path)
    
    elif "sheet" in mime_type or "excel" in mime_type:
        data = process_excel(file_path)

    elif "document" in mime_type or "word" in mime_type:
        data = process_docx(file_path)

    elif any(t in mime_type for t in SimpleTextExtensions):
        data = process_plain_text(file_path)
    
    elif any(t in mime_type for t in ImageExtensions):
        data = process_image(file_path)
    
    
    else:
        raise Exception("This extension is not supported")

    return data